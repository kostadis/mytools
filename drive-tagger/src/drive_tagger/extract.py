"""Turn a Drive file into plain text for embedding.

Native Google editor files are exported as text via gdrive-cli; binary files
(PDF/docx/pptx/txt/md/csv/json) are read either from the optional locally
mounted Drive (matched by md5, to avoid a download) or via ``gdrive-cli cat``.
"""

from __future__ import annotations

import hashlib
import io
import json
import os
import sys
from collections import defaultdict
from functools import lru_cache
from pathlib import Path
from typing import Optional

from . import gdrive
from .config import CONFIG

GOOGLE_APP_PREFIX = "application/vnd.google-apps."

# Native Google editor types we can export, mapped to an export MIME type.
GOOGLE_EXPORT_MIME = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.presentation": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
}

# Native Google types that carry no extractable text body.
GOOGLE_SKIP = {
    "application/vnd.google-apps.folder",
    "application/vnd.google-apps.shortcut",
    "application/vnd.google-apps.form",
    "application/vnd.google-apps.map",
    "application/vnd.google-apps.site",
    "application/vnd.google-apps.drawing",
    "application/vnd.google-apps.fusiontable",
}

TEXT_MIMES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/tab-separated-values",
    "application/json",
    "application/xml",
    "text/xml",
}

PDF_MIME = "application/pdf"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".xml", ".log", ".rst"}

# "Real documents" set used when CONFIG.processable_mode == "documents":
# PDFs, Word/PowerPoint, native Google Docs/Slides, and markdown. Excludes
# text/plain, JSON, CSV, XML, and Google Sheets, which tend to be noise.
DOCUMENT_MIMES = {
    PDF_MIME,
    DOCX_MIME,
    PPTX_MIME,
    "application/vnd.google-apps.document",
    "application/vnd.google-apps.presentation",
    "text/markdown",
}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".markdown"}


def folder_descendants(folder_id: str, scan_path: Path) -> set[str]:
    """Return all file/folder IDs that are descendants of folder_id.

    Reads scan_path once to build a parent→children map, then does a BFS
    from folder_id.  Includes the folder itself so callers can check
    ``rec['id'] in descendants`` or ``any(p in descendants for p in parents)``.
    """
    children: dict[str, list[str]] = defaultdict(list)
    with open(scan_path, "r", encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                rec = json.loads(line)
            except json.JSONDecodeError:
                continue
            for parent in rec.get("parents") or []:
                children[parent].append(rec["id"])

    result: set[str] = {folder_id}
    queue = [folder_id]
    while queue:
        node = queue.pop()
        for child in children.get(node, []):
            if child not in result:
                result.add(child)
                queue.append(child)
    return result


def load_worklist(folder_id: Optional[str], scan_path: Path) -> list[dict]:
    """Return processable Drive records from scan_path in a single file pass.

    If folder_id is given, restricts to descendants of that folder (BFS done
    from the same pass).  This replaces the previous two-pass pattern of
    calling folder_descendants() then reading the file again.
    """
    if not scan_path.exists():
        return []

    children: dict[str, list[str]] = defaultdict(list)
    all_records: list[dict] = []

    with open(scan_path, "r", encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                rec = json.loads(line)
            except json.JSONDecodeError:
                continue
            all_records.append(rec)
            if folder_id:
                for parent in rec.get("parents") or []:
                    children[parent].append(rec["id"])

    # opt-in rpg-lib filename filter
    rpg_names: set[str] | None = None
    if CONFIG.rpg_lib_url:
        from . import rpg_lib as _rpg_lib
        rpg_names = _rpg_lib.load_filenames(CONFIG.rpg_lib_url)

    def _rpg_ok(r: dict) -> bool:
        return rpg_names is None or r.get("name", "") in rpg_names

    if folder_id:
        descendants: set[str] = {folder_id}
        queue = [folder_id]
        while queue:
            node = queue.pop()
            for child in children.get(node, []):
                if child not in descendants:
                    descendants.add(child)
                    queue.append(child)
        filtered = [
            r for r in all_records
            if is_processable(r) and _rpg_ok(r)
            and any(p in descendants for p in (r.get("parents") or []))
        ]
    else:
        filtered = [r for r in all_records if is_processable(r) and _rpg_ok(r)]

    if rpg_names is not None:
        name_counts: dict[str, int] = defaultdict(int)
        for r in filtered:
            name_counts[r.get("name", "")] += 1
        collisions = sum(1 for c in name_counts.values() if c > 1)
        if collisions:
            print(
                f"[worklist] {collisions} filename(s) appear in multiple Drive records (all included)",
                file=sys.stderr,
            )
        print(f"[worklist] rpg-lib filter: {len(filtered)} files", file=sys.stderr)

    return _apply_shard(filtered)


def _apply_shard(records: list[dict]) -> list[dict]:
    """Filter records to the shard specified by DT_SHARD='index/total'."""
    spec = os.environ.get("DT_SHARD", "").strip()
    if not spec:
        return records
    try:
        idx_str, total_str = spec.split("/")
        idx, total = int(idx_str), int(total_str)
        if total < 2 or not (0 <= idx < total):
            raise ValueError
    except (ValueError, TypeError):
        raise ValueError(f"DT_SHARD must be 'index/total' (e.g. '0/2'), got {spec!r}")
    return [
        r for r in records
        if int(hashlib.md5(r["id"].encode()).hexdigest(), 16) % total == idx
    ]


def is_processable(record: dict) -> bool:
    """True if we can plausibly extract text from this record.

    Honors CONFIG.processable_mode: "documents" (default) keeps only real
    documents; "all" also accepts text/plain, JSON, CSV, XML, and Sheets.
    """
    mime = record.get("mime_type", "")
    name = record.get("name", "")
    if record.get("trashed"):
        return False
    if mime in GOOGLE_SKIP:
        return False

    if CONFIG.processable_mode == "documents":
        if mime in DOCUMENT_MIMES:
            return True
        if mime.startswith(GOOGLE_APP_PREFIX):
            return False
        return Path(name).suffix.lower() in DOCUMENT_EXTENSIONS

    # "all" mode
    if mime in GOOGLE_EXPORT_MIME:
        return True
    if mime in (PDF_MIME, DOCX_MIME, PPTX_MIME):
        return True
    if mime in TEXT_MIMES:
        return True
    if mime.startswith(GOOGLE_APP_PREFIX):
        return False
    # Fall back to extension for octet-stream / unknown types.
    return Path(name).suffix.lower() in TEXT_EXTENSIONS


# --- optional local-mount fast path -----------------------------------------


@lru_cache(maxsize=1)
def _md5_index() -> dict[str, str]:
    """Build a {md5: path} index of the mounted Drive (opt-in via DT_USE_MOUNT=1).

    Streaming-mode Drive placeholders are skipped (their md5 won't match). The
    index is built once per process and cached.
    """
    if os.environ.get("DT_USE_MOUNT", "0") != "1":
        return {}
    root = Path(CONFIG.drive_mount)
    if not root.is_dir():
        return {}
    index: dict[str, str] = {}
    for dirpath, _dirs, files in os.walk(root):
        for fn in files:
            p = Path(dirpath) / fn
            try:
                if not p.is_file() or p.stat().st_size == 0:
                    continue
                h = hashlib.md5()
                with open(p, "rb") as fh:
                    for chunk in iter(lambda: fh.read(1 << 20), b""):
                        h.update(chunk)
                index.setdefault(h.hexdigest(), str(p))
            except OSError:
                continue
    return index


def _local_bytes(md5: Optional[str]) -> Optional[bytes]:
    if not md5:
        return None
    path = _md5_index().get(md5)
    if not path:
        return None
    try:
        return Path(path).read_bytes()
    except OSError:
        return None


# --- binary parsers ----------------------------------------------------------


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx_text(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append("".join(run.text for run in para.runs))
    return "\n".join(p for p in parts if p)


def _decode_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def extract_text(record: dict) -> Optional[str]:
    """Extract plain text for a Drive record, or ``None`` if not extractable.

    The result is truncated to ``CONFIG.max_chars`` for embedding stability.
    """
    if not is_processable(record):
        return None

    file_id = record["id"]
    mime = record.get("mime_type", "")
    name = record.get("name", "")

    text: Optional[str] = None
    if mime in GOOGLE_EXPORT_MIME:
        text = gdrive.export(file_id, mime=GOOGLE_EXPORT_MIME[mime])
    else:
        data = _local_bytes(record.get("md5_checksum")) or gdrive.cat(file_id)
        if mime == PDF_MIME:
            text = _pdf_text(data)
        elif mime == DOCX_MIME:
            text = _docx_text(data)
        elif mime == PPTX_MIME:
            text = _pptx_text(data)
        elif mime in TEXT_MIMES or Path(name).suffix.lower() in TEXT_EXTENSIONS:
            text = _decode_text(data)

    if text is None:
        return None
    text = text.strip()
    if not text:
        return None
    return text[: CONFIG.max_chars]
